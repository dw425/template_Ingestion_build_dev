import os
import json
import re
from datetime import datetime
from typing import Dict, List, Any, Optional
import PyPDF2
from docx import Document
from pptx import Presentation
import openai
import anthropic
import tiktoken

class AIDocumentAnalyzer:
    def __init__(self):
        self.openai_client = None
        self.anthropic_client = None
        self.setup_ai_clients()
        
    def setup_ai_clients(self):
        """Setup AI clients based on available API keys"""
        openai_key = os.getenv('OPENAI_API_KEY')
        anthropic_key = os.getenv('ANTHROPIC_API_KEY')
        
        if openai_key:
            openai.api_key = openai_key
            self.openai_client = openai
        
        if anthropic_key:
            self.anthropic_client = anthropic.Anthropic(api_key=anthropic_key)
    
    def analyze_document(self, filepath: str) -> Dict[str, Any]:
        """Main entry point for document analysis"""
        try:
            # Extract text from document
            text_content = self.extract_text_from_file(filepath)
            
            if not text_content.strip():
                raise Exception("No text content could be extracted from the document")
            
            # Use AI to analyze the content
            analysis_results = self.ai_analyze_content(text_content, filepath)
            
            return {
                'file_path': filepath,
                'text_length': len(text_content),
                'ai_analysis': analysis_results,
                'extracted_text_preview': text_content[:500] + "..." if len(text_content) > 500 else text_content
            }
            
        except Exception as e:
            raise Exception(f"Document analysis failed: {str(e)}")
    
    def extract_text_from_file(self, filepath: str) -> str:
        """Extract text from various file formats"""
        file_extension = filepath.lower().split('.')[-1]
        
        if file_extension == 'pdf':
            return self.extract_from_pdf(filepath)
        elif file_extension == 'docx':
            return self.extract_from_docx(filepath)
        elif file_extension in ['pptx', 'ppt']:
            return self.extract_from_pptx(filepath)
        else:
            raise Exception(f"Unsupported file type: {file_extension}")
    
    def extract_from_pdf(self, filepath: str) -> str:
        """Extract text from PDF files"""
        text = ""
        try:
            with open(filepath, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            raise Exception(f"PDF extraction failed: {str(e)}")
        return text
    
    def extract_from_docx(self, filepath: str) -> str:
        """Extract text from DOCX files"""
        try:
            doc = Document(filepath)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
            
            return text
        except Exception as e:
            raise Exception(f"DOCX extraction failed: {str(e)}")
    
    def extract_from_pptx(self, filepath: str) -> str:
        """Extract text from PPTX files"""
        try:
            presentation = Presentation(filepath)
            text = ""
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                text += f"\n--- Slide {slide_num} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                    
                    # Extract text from tables in slides
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                text += cell.text + " "
                            text += "\n"
            
            return text
        except Exception as e:
            raise Exception(f"PPTX extraction failed: {str(e)}")
    
    def ai_analyze_content(self, text_content: str, filepath: str) -> Dict[str, Any]:
        """Use AI to analyze document content for project insights"""
        
        # Create analysis prompt
        analysis_prompt = self.create_analysis_prompt(text_content, filepath)
        
        # Try Claude first, then OpenAI as fallback
        if self.anthropic_client:
            return self.analyze_with_claude(analysis_prompt)
        elif self.openai_client:
            return self.analyze_with_openai(analysis_prompt)
        else:
            # Fallback to rule-based analysis
            return self.rule_based_analysis(text_content)
    
    def create_analysis_prompt(self, text_content: str, filepath: str) -> str:
        """Create a comprehensive analysis prompt"""
        return f"""
You are an expert project management analyst. Analyze the following document content and extract key project insights.

Document: {os.path.basename(filepath)}

CONTENT TO ANALYZE:
{text_content}

Please provide a comprehensive analysis in the following JSON format:

{{
    "project_overview": {{
        "project_name": "extracted project name or 'Unknown'",
        "project_type": "type of project (software, construction, marketing, etc.)",
        "description": "brief project description",
        "objectives": ["list of main objectives found"]
    }},
    "timeline_info": {{
        "start_date": "YYYY-MM-DD or null",
        "end_date": "YYYY-MM-DD or null", 
        "milestones": ["list of key milestones with dates if available"],
        "deadlines": ["list of important deadlines"]
    }},
    "team_and_resources": {{
        "team_members": ["list of people mentioned"],
        "roles": ["list of roles/positions mentioned"],
        "departments": ["list of departments involved"],
        "budget_info": "any budget/cost information found"
    }},
    "kpis_and_metrics": {{
        "performance_metrics": ["list of KPIs, metrics, or measurements"],
        "success_criteria": ["list of success criteria"],
        "progress_indicators": ["list of progress indicators"],
        "numerical_targets": ["specific numbers/percentages/targets found"]
    }},
    "status_and_progress": {{
        "current_status": "overall project status",
        "completed_tasks": ["list of completed items"],
        "in_progress_tasks": ["list of ongoing items"],
        "pending_tasks": ["list of pending/planned items"],
        "risks_issues": ["list of risks, issues, or blockers"]
    }},
    "key_insights": {{
        "strengths": ["positive aspects identified"],
        "concerns": ["areas of concern or risk"],
        "recommendations": ["suggested actions or improvements"],
        "critical_success_factors": ["most important factors for success"]
    }},
    "document_metadata": {{
        "document_type": "type of document (report, presentation, plan, etc.)",
        "confidence_score": 0.85,
        "data_quality": "assessment of how much useful data was extracted",
        "missing_information": ["key information that seems to be missing"]
    }}
}}

IMPORTANT: 
- Extract only information that is clearly present in the document
- Use null for dates that cannot be determined
- Be specific and factual, avoid assumptions
- If information is unclear, note it in the confidence assessment
- Focus on actionable project management insights
"""
    
    def analyze_with_claude(self, prompt: str) -> Dict[str, Any]:
        """Analyze using Anthropic Claude"""
        try:
            response = self.anthropic_client.messages.create(
                model="claude-3-sonnet-20240229",
                max_tokens=4000,
                messages=[
                    {
                        "role": "user",
                        "content": prompt
                    }
                ]
            )
            
            # Extract JSON from response
            content = response.content[0].text
            return self.extract_json_from_response(content)
            
        except Exception as e:
            print(f"Claude analysis failed: {str(e)}")
            return self.create_fallback_analysis("Claude API error")
    
    def analyze_with_openai(self, prompt: str) -> Dict[str, Any]:
        """Analyze using OpenAI GPT"""
        try:
            response = self.openai_client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert project management analyst. Always respond with valid JSON."
                    },
                    {
                        "role": "user", 
                        "content": prompt
                    }
                ],
                max_tokens=4000,
                temperature=0.3
            )
            
            content = response.choices[0].message.content
            return self.extract_json_from_response(content)
            
        except Exception as e:
            print(f"OpenAI analysis failed: {str(e)}")
            return self.create_fallback_analysis("OpenAI API error")
    
    def extract_json_from_response(self, content: str) -> Dict[str, Any]:
        """Extract JSON from AI response"""
        try:
            # Try to find JSON in the response
            json_start = content.find('{')
            json_end = content.rfind('}') + 1
            
            if json_start >= 0 and json_end > json_start:
                json_str = content[json_start:json_end]
                return json.loads(json_str)
            else:
                raise Exception("No valid JSON found in response")
                
        except Exception as e:
            print(f"JSON extraction failed: {str(e)}")
            return self.create_fallback_analysis("JSON parsing error")
    
    def rule_based_analysis(self, text_content: str) -> Dict[str, Any]:
        """Fallback rule-based analysis when AI is not available"""
        
        # Basic keyword extraction
        project_keywords = self.extract_project_keywords(text_content)
        dates = self.extract_dates(text_content)
        metrics = self.extract_metrics(text_content)
        
        return {
            "project_overview": {
                "project_name": "Extracted from document",
                "project_type": "Unknown",
                "description": f"Document contains {len(text_content.split())} words",
                "objectives": project_keywords.get('objectives', [])
            },
            "timeline_info": {
                "start_date": dates[0] if dates else None,
                "end_date": dates[-1] if len(dates) > 1 else None,
                "milestones": dates[:5],  # First 5 dates as milestones
                "deadlines": []
            },
            "kpis_and_metrics": {
                "performance_metrics": metrics,
                "numerical_targets": self.extract_numbers(text_content)
            },
            "status_and_progress": {
                "current_status": "Analysis based on keyword extraction",
                "completed_tasks": project_keywords.get('completed', []),
                "in_progress_tasks": project_keywords.get('in_progress', []),
                "pending_tasks": project_keywords.get('pending', [])
            },
            "document_metadata": {
                "document_type": "Unknown document type",
                "confidence_score": 0.3,
                "data_quality": "Limited - rule-based extraction only",
                "analysis_method": "Keyword extraction (AI not available)"
            }
        }
    
    def extract_project_keywords(self, text: str) -> Dict[str, List[str]]:
        """Extract project-related keywords"""
        completed_patterns = [
            r'completed?\s+([^.]+)',
            r'finished\s+([^.]+)',
            r'done\s+([^.]+)'
        ]
        
        in_progress_patterns = [
            r'working on\s+([^.]+)',
            r'in progress\s+([^.]+)',
            r'currently\s+([^.]+)'
        ]
        
        pending_patterns = [
            r'planned?\s+([^.]+)',
            r'upcoming\s+([^.]+)',
            r'next\s+([^.]+)'
        ]
        
        return {
            'completed': self.extract_with_patterns(text, completed_patterns),
            'in_progress': self.extract_with_patterns(text, in_progress_patterns),
            'pending': self.extract_with_patterns(text, pending_patterns)
        }
    
    def extract_with_patterns(self, text: str, patterns: List[str]) -> List[str]:
        """Extract text using regex patterns"""
        results = []
        for pattern in patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            results.extend([match.strip() for match in matches])
        return results[:5]  # Limit to 5 results
    
    def extract_dates(self, text: str) -> List[str]:
        """Extract dates from text"""
        date_patterns = [
            r'\d{4}-\d{2}-\d{2}',
            r'\d{1,2}/\d{1,2}/\d{4}',
            r'(January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{1,2},?\s+\d{4}'
        ]
        
        dates = []
        for pattern in date_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            dates.extend(matches)
        
        return list(set(dates))[:10]  # Unique dates, max 10
    
    def extract_metrics(self, text: str) -> List[str]:
        """Extract potential metrics and KPIs"""
        metric_keywords = [
            'revenue', 'profit', 'cost', 'budget', 'roi', 'performance',
            'efficiency', 'productivity', 'quality', 'satisfaction',
            'completion rate', 'progress', 'milestone', 'target'
        ]
        
        metrics = []
        for keyword in metric_keywords:
            pattern = rf'{keyword}[:\s]*([^.]+)'
            matches = re.findall(pattern, text, re.IGNORECASE)
            metrics.extend([f"{keyword}: {match.strip()}" for match in matches])
        
        return metrics[:10]  # Max 10 metrics
    
    def extract_numbers(self, text: str) -> List[str]:
        """Extract numbers that might be targets or metrics"""
        number_patterns = [
            r'\d+%',  # Percentages
            r'\$[\d,]+',  # Currency
            r'\d+\s*(days?|weeks?|months?)',  # Time periods
            r'\d+\.\d+',  # Decimals
        ]
        
        numbers = []
        for pattern in number_patterns:
            matches = re.findall(pattern, text)
            numbers.extend(matches)
        
        return list(set(numbers))[:15]  # Unique numbers, max 15
    
    def create_fallback_analysis(self, error_reason: str) -> Dict[str, Any]:
        """Create a basic analysis structure when AI fails"""
        return {
            "project_overview": {
                "project_name": "Unknown",
                "project_type": "Unknown", 
                "description": "Could not analyze document content",
                "objectives": []
            },
            "timeline_info": {
                "start_date": None,
                "end_date": None,
                "milestones": [],
                "deadlines": []
            },
            "team_and_resources": {
                "team_members": [],
                "roles": [],
                "departments": [],
                "budget_info": "Not available"
            },
            "kpis_and_metrics": {
                "performance_metrics": [],
                "success_criteria": [],
                "progress_indicators": [],
                "numerical_targets": []
            },
            "status_and_progress": {
                "current_status": "Analysis unavailable",
                "completed_tasks": [],
                "in_progress_tasks": [],
                "pending_tasks": [],
                "risks_issues": []
            },
            "key_insights": {
                "strengths": [],
                "concerns": [f"Analysis failed: {error_reason}"],
                "recommendations": ["Upload a supported file format for AI analysis"],
                "critical_success_factors": []
            },
            "document_metadata": {
                "document_type": "Unknown",
                "confidence_score": 0.0,
                "data_quality": "Analysis failed",
                "missing_information": ["All project information - analysis unavailable"],
                "error": error_reason
            }
        }
